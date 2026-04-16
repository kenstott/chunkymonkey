# Copyright (c) 2025 Kenneth Stott. MIT License.
# Canary: 32dc2e60-6d2c-4800-b55b-47734ece8b37
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Generate demo office documents for chunkymonkey — xlsx, docx, pptx.

These documents are designed to demonstrate that contextual chunking
(section breadcrumbs in embedding_content) outperforms naive chunking
on realistic office document structures.

Each document is structured to create the same "context loss" problem
as the Markdown demos:
  - Repeated section/sheet names across regions or categories
  - Long tables that split across chunk boundaries
  - Parent-level context (region, category) only in the heading, not in rows

Run::
    cd /path/to/chunkymonkey
    pip install openpyxl python-docx python-pptx
    python demo/generate_office_docs.py
"""
from __future__ import annotations

from pathlib import Path

HERE = Path(__file__).parent
DOCS_DIR = HERE / "docs"
DOCS_DIR.mkdir(exist_ok=True)


# ─────────────────────────────────────────────────────────────────────────────
# Excel — Regional Sales Pipeline (xlsx)
# ─────────────────────────────────────────────────────────────────────────────

def generate_xlsx() -> Path:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        print("  SKIP xlsx — pip install openpyxl")
        return DOCS_DIR / "sales_pipeline.xlsx"

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # Remove default sheet

    REGIONS = ["Americas", "EMEA", "APAC"]
    PRODUCTS = [
        ("Shield Elite", "XDR", "Enterprise"),
        ("TrustBroker Advanced", "Zero Trust", "Enterprise"),
        ("KeyVault Enterprise", "PAM", "Enterprise"),
        ("Sentinel Pro", "EDR", "Mid-Market"),
        ("AccessGuard Premium", "PAM", "Mid-Market"),
        ("Guardian Premium", "XDR", "Mid-Market"),
        ("CredShield Pro", "MFA", "SMB"),
        ("Bastion Agent Pro", "EDR", "SMB"),
        ("DocVault Pro", "Versioning", "SMB"),
        ("Rampart Standard", "EDR", "SMB"),
        ("Fortress XDR", "XDR", "Enterprise"),
        ("Aegis Endpoint+", "XDR", "Enterprise"),
        ("TrustBroker", "Zero Trust", "Mid-Market"),
        ("KeyVault Pro", "MFA+SSO", "Mid-Market"),
        ("AccessGuard Standard", "MFA+SSO", "SMB"),
        ("Bulwark Endpoint+", "XDR", "Enterprise"),
        ("Barricade Elite", "XDR", "Enterprise"),
        ("Citadel Pro", "XDR", "Mid-Market"),
        ("Vanguard Endpoint Elite", "XDR", "Enterprise"),
        ("Parapet Advanced", "EDR", "Mid-Market"),
    ]

    STAGES = ["Prospecting", "Discovery", "Proposal", "Negotiation", "Closed Won"]

    import random
    rng = random.Random(42)

    header_fill = PatternFill("solid", fgColor="2E75B6")
    header_font = Font(bold=True, color="FFFFFF")

    for i, region in enumerate(REGIONS):
        ws = wb.create_sheet(title=region)

        # Title row — region context (mirrors the markdown heading problem:
        # this only appears once, at the top of the sheet)
        ws["A1"] = f"{region} — Sales Pipeline Q1 2024"
        ws["A1"].font = Font(bold=True, size=14)
        ws.merge_cells("A1:G1")

        # Header row
        headers = ["Deal ID", "Account", "Product", "Category", "Segment", "Stage", "ARR ($k)"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")

        # Data rows — 25 deals per region
        # NOTE: rows contain product/stage/ARR but NOT the region name.
        # Region context lives only in the sheet title (row 1) and sheet tab.
        # This is the "context loss" point for chunking.
        deal_num = i * 100 + 1
        for row_idx in range(25):
            product, category, segment = PRODUCTS[row_idx % len(PRODUCTS)]
            stage = STAGES[rng.randint(0, len(STAGES) - 1)]
            arr = rng.randint(50, 800)
            account = f"Acct-{rng.randint(1000, 9999)}"
            ws.cell(row=4 + row_idx, column=1, value=f"{region[:2].upper()}-{deal_num + row_idx:04d}")
            ws.cell(row=4 + row_idx, column=2, value=account)
            ws.cell(row=4 + row_idx, column=3, value=product)
            ws.cell(row=4 + row_idx, column=4, value=category)
            ws.cell(row=4 + row_idx, column=5, value=segment)
            ws.cell(row=4 + row_idx, column=6, value=stage)
            ws.cell(row=4 + row_idx, column=7, value=arr)

        # Column widths
        ws.column_dimensions["A"].width = 14
        ws.column_dimensions["B"].width = 14
        ws.column_dimensions["C"].width = 26
        ws.column_dimensions["D"].width = 14
        ws.column_dimensions["E"].width = 14
        ws.column_dimensions["F"].width = 16
        ws.column_dimensions["G"].width = 10

    out = DOCS_DIR / "sales_pipeline.xlsx"
    wb.save(out)
    print(f"  Generated {out.name}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# Word — Engineering Architecture Review (docx)
# ─────────────────────────────────────────────────────────────────────────────

def generate_docx() -> Path:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        print("  SKIP docx — pip install python-docx")
        return DOCS_DIR / "architecture_review.docx"

    doc = Document()
    doc.add_heading("Engineering Architecture Review — FY2024", level=0)
    doc.add_paragraph(
        "This document covers architecture decisions, component inventories, "
        "and risk assessments for all three regional engineering hubs. "
        "Each region runs an independent deployment topology but shares the "
        "same core platform services."
    )

    REGIONS = [
        ("Americas", "US-EAST-1 / US-WEST-2", "Kubernetes 1.28", "PostgreSQL 15", "Redis 7.2"),
        ("EMEA",     "EU-WEST-1 / EU-CENTRAL-1", "Kubernetes 1.27", "PostgreSQL 14", "Redis 7.0"),
        ("APAC",     "AP-SOUTHEAST-1 / AP-NORTHEAST-1", "Kubernetes 1.28", "PostgreSQL 15", "Redis 7.2"),
    ]

    COMPONENTS = [
        ("Auth Service",       "Go 1.21",      "OAuth 2.0 / OIDC",      "Active-active"),
        ("API Gateway",        "Envoy 1.28",   "REST / gRPC",           "Active-active"),
        ("Data Pipeline",      "Python 3.12",  "Kafka 3.6",             "Active-passive"),
        ("Search Service",     "Java 21",      "Elasticsearch 8.10",    "Active-active"),
        ("Notification Queue", "Go 1.21",      "RabbitMQ 3.12",         "Active-passive"),
        ("Object Storage",     "Rust 1.74",    "S3-compatible API",     "Active-active"),
        ("Scheduler",          "Python 3.12",  "Celery 5.3",            "Active-passive"),
        ("Database Primary",   "PostgreSQL",   "pgBouncer 1.21",        "Primary-replica"),
        ("CDN Edge",           "Nginx 1.25",   "PoP mesh",              "Active-active"),
        ("Monitoring Stack",   "Go 1.21",      "Prometheus / Grafana",  "Active-passive"),
        ("Identity Provider",  "Keycloak 22",  "SAML 2.0 / LDAP",       "Active-passive"),
        ("Feature Flags",      "Go 1.21",      "gRPC",                  "Active-active"),
        ("Rate Limiter",       "C++ / WASM",   "Redis-backed",          "Active-active"),
        ("Audit Log",          "Rust 1.74",    "Kafka → S3",            "Write-once"),
        ("Config Service",     "Go 1.21",      "etcd 3.5",              "Active-passive"),
    ]

    RISKS = [
        ("CDN single-PoP dependency", "High", "Multi-PoP failover — Q2"),
        ("Auth SPOF during key rotation", "Medium", "Blue-green rotation — Q1"),
        ("Search OOM on index rebuild", "High", "Memory limits + incremental rebuild"),
        ("Kafka consumer lag on peak", "Medium", "Auto-scaling consumer groups"),
        ("DB autovacuum contention", "Low", "Scheduled maintenance windows"),
    ]

    for region, az, k8s, db, cache in REGIONS:
        doc.add_heading(region, level=1)
        doc.add_paragraph(
            f"Primary availability zones: {az}. "
            f"Container orchestration: {k8s}. "
            f"Primary datastore: {db}. "
            f"Cache layer: {cache}."
        )

        # Component inventory — long table that WILL split across chunks.
        # Continuation rows contain only component name/tech/protocol/topology.
        # The region heading is NOT in those rows.
        doc.add_heading("Component Inventory", level=2)
        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "Component"
        hdr[1].text = "Technology"
        hdr[2].text = "Protocol"
        hdr[3].text = "Topology"

        import random
        rng = random.Random(hash(region))
        for name, tech, proto, topo in COMPONENTS:
            row = table.add_row().cells
            row[0].text = name
            row[1].text = tech
            row[2].text = proto
            row[3].text = topo

        doc.add_paragraph("")  # spacing

        doc.add_heading("Risk Register", level=2)
        risk_table = doc.add_table(rows=1, cols=3)
        risk_table.style = "Light Grid Accent 1"
        rh = risk_table.rows[0].cells
        rh[0].text = "Risk"
        rh[1].text = "Severity"
        rh[2].text = "Mitigation"
        for risk, sev, mit in RISKS:
            rrow = risk_table.add_row().cells
            rrow[0].text = risk
            rrow[1].text = sev
            rrow[2].text = mit

        doc.add_paragraph("")

    out = DOCS_DIR / "architecture_review.docx"
    doc.save(out)
    print(f"  Generated {out.name}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint — Quarterly Business Review (pptx)
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx() -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print("  SKIP pptx — pip install python-pptx")
        return DOCS_DIR / "qbr_presentation.pptx"

    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # Title and Content

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Q1 2024 Quarterly Business Review"
    title_slide.placeholders[1].text = "All Regions — Confidential"

    REGIONS = ["Americas", "EMEA", "APAC"]

    METRICS = [
        # (metric, Americas, EMEA, APAC)
        ("New Logos",            "142",  "98",  "61"),
        ("Expansion ARR ($M)",   "$4.2", "$2.9", "$1.8"),
        ("Churn Rate",           "2.1%", "1.8%", "3.2%"),
        ("NPS Score",            "52",   "48",   "44"),
        ("Support CSAT",         "4.4",  "4.2",  "4.1"),
        ("P0 Incidents",         "8",    "6",    "4"),
        ("Mean MTTR (min)",      "58",   "64",   "72"),
        ("Deployment Frequency", "4.2/d","3.8/d","3.1/d"),
        ("Lead Time (days)",     "1.4",  "1.6",  "2.1"),
        ("Change Fail Rate",     "3.2%", "4.1%", "5.6%"),
    ]

    # Executive summary slide
    exec_slide = prs.slides.add_slide(blank_layout)
    exec_slide.shapes.title.text = "Executive Summary"
    tf = exec_slide.placeholders[1].text_frame
    tf.text = "Q1 2024 highlights across all regions:"
    tf.add_paragraph().text = "• Americas: strong new logo growth, stable churn"
    tf.add_paragraph().text = "• EMEA: best NPS and churn; P0 reduction YoY"
    tf.add_paragraph().text = "• APAC: emerging market growth; MTTR focus area"

    for region in REGIONS:
        # Regional overview slide — region name IS here (heading)
        overview = prs.slides.add_slide(blank_layout)
        overview.shapes.title.text = f"{region} — Overview"
        tf = overview.placeholders[1].text_frame
        tf.text = f"Q1 2024 performance summary for the {region} region."
        tf.add_paragraph().text = ""

        # Key metrics slide — metrics table, region name ONLY in title
        metrics_slide = prs.slides.add_slide(blank_layout)
        metrics_slide.shapes.title.text = f"{region} — Key Metrics"
        tf2 = metrics_slide.placeholders[1].text_frame
        # Simulate a metrics table as bullet points
        # (pptx tables require more complex layout; bullets are cleaner for demo)
        col_idx = REGIONS.index(region)
        tf2.text = "Metric  /  Value"
        for metric, am, emea, apac in METRICS:
            value = [am, emea, apac][col_idx]
            tf2.add_paragraph().text = f"{metric}: {value}"

        # Incidents slide — long bullet list, region only in title
        inc_slide = prs.slides.add_slide(blank_layout)
        inc_slide.shapes.title.text = f"{region} — Incident Review"
        tf3 = inc_slide.placeholders[1].text_frame
        tf3.text = "P0 Incidents Q1 2024:"
        incidents = [
            ("Auth service timeout", "Jan", "94 min"),
            ("CDN edge BGP flap",    "Feb", "68 min"),
            ("Search OOM rebuild",   "Mar", "112 min"),
            ("DB write replica full","Jan", "47 min"),
            ("API gateway config",   "Feb", "31 min"),
        ]
        for svc, month, mttr in incidents:
            tf3.add_paragraph().text = f"{month} — {svc} (MTTR: {mttr})"

    out = DOCS_DIR / "qbr_presentation.pptx"
    prs.save(out)
    print(f"  Generated {out.name}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating demo office documents...")
    generate_xlsx()
    generate_docx()
    generate_pptx()
    print("\nDone. Files written to demo/docs/")
    print(
        "\nThese documents share the same structural property as the Markdown demos:\n"
        "  • Parent section / sheet / slide title contains the REGION name\n"
        "  • Row-level data contains product/metric/incident rows WITHOUT the region\n"
        "  • When chunked naively, continuation chunks are indistinguishable across regions\n"
        "  • Contextual chunking (section breadcrumb in embedding_content) resolves this\n"
    )
