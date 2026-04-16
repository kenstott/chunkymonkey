# Copyright (c) 2025 Kenneth Stott. MIT License.
# Canary: fb48752f-c49b-427f-a135-6c6a75f90ed4
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""PPTX text extractor using python-pptx."""

from __future__ import annotations

from io import BytesIO

try:
    import pptx as _pptx_module
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


def _extract_pptx_content(prs) -> str:
    """Extract text content from a python-pptx Presentation object."""
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            if shape.has_table:
                table_rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    slide_text.append("\n".join(table_rows))

        if slide_text:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))

    return "\n\n".join(slides)


class PptxExtractor:
    """Extract plain text from PPTX bytes."""

    def can_handle(self, doc_type: str) -> bool:
        return doc_type == "pptx"

    def extract(self, data: bytes, source_path: str | None = None) -> str:
        if not _PPTX_AVAILABLE:
            raise ImportError("pip install chunkymonkey[pptx]")

        prs = _pptx_module.Presentation(BytesIO(data))
        return _extract_pptx_content(prs)
